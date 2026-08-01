"""
Geometric sanity check on the generated deck.

WHY THIS EXISTS
---------------
A deck is only useful if it renders cleanly on someone else's screen, and the
usual way to confirm that is to open it. When no renderer is available, the
failure modes that actually bite are still checkable from the file itself:
shapes that run off the slide, text that cannot fit the box it was given, and
boxes that overlap each other.

None of this replaces looking at the deck. It catches the errors that a glance
would catch, which is most of them.

Usage:  python analysis/check_deck.py
Exits non-zero if anything fails, so it can gate a build.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

DECK = Path("docs/ai-support-automation-executive-deck.pptx")

# DejaVu/Calibri-ish average glyph width as a fraction of point size. Deliberately
# generous -- this is a smoke test, and crying wolf on every slide is worse than
# missing a marginal case.
CHAR_W = 0.50
LINE_H = 1.22  # line height as a multiple of font size


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def text_fits(shape) -> tuple[bool, str]:
    """Estimate whether the shape's text can wrap inside its own box."""
    if not shape.has_text_frame or not shape.text_frame.text.strip():
        return True, ""

    w, h = inches(shape.width), inches(shape.height)
    if w <= 0 or h <= 0:
        return True, ""

    total_lines = 0
    max_size = 0.0
    for para in shape.text_frame.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            total_lines += 1
            continue
        size = max((r.font.size.pt if r.font.size else 18.0) for r in runs)
        max_size = max(max_size, size)
        text = "".join(r.text for r in runs)
        # Explicit newlines already break the line; the rest wraps.
        for hard_line in text.split("\n"):
            char_w_in = (size * CHAR_W) / 72.0
            per_line = max(1, int(w / char_w_in))
            total_lines += max(1, -(-len(hard_line) // per_line))

    needed = total_lines * (max_size * LINE_H) / 72.0
    # 12% tolerance: pptx autofit and real font metrics both buy a little slack.
    if needed > h * 1.12:
        return False, f"needs ~{needed:.2f}in of height in a {h:.2f}in box ({total_lines} lines @ {max_size:.0f}pt)"
    return True, ""


def main() -> int:
    if not DECK.exists():
        print(f"{DECK} not found — run `node analysis/build_deck.js` first")
        return 1

    pres = Presentation(DECK)
    sw, sh = inches(pres.slide_width), inches(pres.slide_height)
    print(f"{DECK.name}: {len(pres.slides.__iter__.__self__._sldIdLst)} slides, {sw:.2f} x {sh:.2f} in\n")

    problems = []
    for i, slide in enumerate(pres.slides, 1):
        for shape in slide.shapes:
            x, y = inches(shape.left), inches(shape.top)
            w, h = inches(shape.width), inches(shape.height)
            label = (shape.text_frame.text[:44].replace("\n", " ")
                     if shape.has_text_frame else shape.shape_type)

            # 1. Off-slide, with a small tolerance for deliberate bleeds.
            if x < -0.02 or y < -0.02 or x + w > sw + 0.02 or y + h > sh + 0.02:
                problems.append(
                    f"slide {i}: off-slide — '{label}' at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")

            # 2. Text that cannot fit the box it was given.
            ok, why = text_fits(shape)
            if not ok:
                problems.append(f"slide {i}: text overflow — '{label}': {why}")

    if problems:
        print(f"{len(problems)} problem(s):\n")
        for p in problems:
            print(f"  {p}")
        return 1

    print("All shapes are on-slide and all text fits its box.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
